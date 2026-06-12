"""Genera el reporte ejecutivo mensual de TI a partir de la plantilla de captura.

Uso:
    python3 generar_reporte.py [Reporte_TI_Plantilla.xlsx]

Salida: Reporte_TI_Ejecutivo_<Mes>_<Anio>.pptx
"""
import sys
from datetime import datetime, date
from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ============= Paleta corporativa =============
GREEN_DEEP = RGBColor(0x1F, 0x4E, 0x2C)
GREEN_MID = RGBColor(0x3E, 0x7C, 0x47)
GREEN_LIGHT = RGBColor(0xD9, 0xE8, 0xDA)
AMBER = RGBColor(0xD4, 0xA2, 0x4C)
AMBER_LIGHT = RGBColor(0xF4, 0xE5, 0xC2)
CREAM = RGBColor(0xF8, 0xF4, 0xE9)
DARK = RGBColor(0x1A, 0x2E, 0x1A)
GRAY = RGBColor(0x5A, 0x5A, 0x5A)
GRAY_LIGHT = RGBColor(0xE5, 0xE5, 0xE5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x39, 0x2B)
YELLOW = RGBColor(0xE6, 0xB3, 0x3B)
GREEN_OK = RGBColor(0x3C, 0x8F, 0x4F)

HEAD_FONT = "Calibri"
BODY_FONT = "Calibri"

# ============= Utilidades =============

def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, *, font=BODY_FONT, size=14, bold=False,
             color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(text, list):
        for i, line in enumerate(text):
            run = p.add_run() if i == 0 else tf.add_paragraph().add_run()
            if i > 0:
                run.font.size = Pt(size)
            run.text = line
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
    else:
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=0):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(line_width)
    sh.shadow.inherit = False
    return sh


def add_circle(slide, x, y, d, fill_color):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


STATE_COLORS = {
    "Verde": GREEN_OK,
    "Amarillo": YELLOW,
    "Rojo": RED,
    "Completado": GREEN_DEEP,
    "Crítica": RED,
    "Alta": YELLOW,
    "Media": AMBER,
    "Baja": GRAY,
}


def state_pill(slide, x, y, w, h, label):
    color = STATE_COLORS.get(label, GRAY)
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = HEAD_FONT
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = WHITE


def page_chrome(slide, page_num, total, mes_anio):
    """Banda inferior con número de página y mes/año."""
    add_rect(slide, 0, 7.05, 13.333, 0.45, GREEN_DEEP)
    add_text(slide, 0.5, 7.10, 6, 0.35, "Reporte Ejecutivo de TI · Ingenio La Corona",
             size=10, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, 7.5, 7.10, 5.3, 0.35, f"{mes_anio}    ·    {page_num} / {total}",
             size=10, color=AMBER_LIGHT, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)


def slide_title(slide, kicker, title):
    add_rect(slide, 0.5, 0.55, 0.10, 0.55, AMBER)
    add_text(slide, 0.75, 0.50, 12, 0.30, kicker.upper(),
             size=11, bold=True, color=AMBER, anchor=MSO_ANCHOR.TOP)
    add_text(slide, 0.75, 0.78, 12, 0.50, title,
             font=HEAD_FONT, size=28, bold=True, color=GREEN_DEEP, anchor=MSO_ANCHOR.TOP)


# ============= Lectura del Excel =============

def read_workbook(path):
    wb = load_workbook(path, data_only=True)
    data = {}

    ws = wb["Inicio"]
    inicio = {}
    inicio["mes"] = ws.cell(row=4, column=2).value or "Mes"
    inicio["anio"] = ws.cell(row=5, column=2).value or datetime.now().year
    inicio["zafra"] = ws.cell(row=6, column=2).value or ""
    inicio["responsable"] = ws.cell(row=7, column=2).value or ""
    fecha = ws.cell(row=8, column=2).value
    if isinstance(fecha, (datetime, date)):
        inicio["fecha"] = fecha.strftime("%Y-%m-%d")
    else:
        inicio["fecha"] = str(fecha) if fecha else ""
    inicio["mensaje"] = ws.cell(row=9, column=2).value or ""
    puntos = []
    for r in range(12, 17):
        v = ws.cell(row=r, column=2).value
        if v:
            puntos.append(str(v))
    inicio["puntos"] = puntos
    data["inicio"] = inicio

    def read_table(sheet, start_row, columns, end_marker_col=1):
        ws = wb[sheet]
        rows = []
        r = start_row
        while True:
            v = ws.cell(row=r, column=end_marker_col).value
            if v is None or str(v).strip() == "":
                break
            row = {}
            for col_idx, key in enumerate(columns, start=1):
                row[key] = ws.cell(row=r, column=col_idx).value
            rows.append(row)
            r += 1
        return rows

    data["disponibilidad"] = read_table(
        "Disponibilidad", 5,
        ["sistema", "categoria", "obj", "real", "horas", "estado", "comentario"])
    data["incidentes"] = read_table(
        "Incidentes", 5,
        ["fecha", "sistema", "severidad", "descripcion", "causa", "horas", "accion"])
    data["proyectos"] = read_table(
        "Proyectos", 5,
        ["proyecto", "responsable", "estado", "avance", "inicio", "fin", "hito", "riesgo"])
    presup = read_table(
        "Presupuesto", 5,
        ["categoria", "presupuesto", "ejecutado", "variacion", "ejec_pct", "comentario"])
    presup = [r for r in presup if str(r["categoria"]).upper() != "TOTAL"]
    data["presupuesto"] = presup
    data["seguridad"] = read_table(
        "Seguridad", 5,
        ["indicador", "actual", "anterior", "tendencia", "comentario"])
    data["licencias"] = read_table(
        "Licencias", 5,
        ["producto", "proveedor", "cantidad", "vence", "dias", "estado", "costo"])
    data["inventario"] = read_table(
        "Inventario", 5,
        ["categoria", "cantidad", "edad", "rep_3m", "rep_12m", "comentario"])

    return data


# ============= Slides =============

def fmt_money(x):
    if x is None or x == "":
        return ""
    try:
        return f"${float(x):,.0f}"
    except Exception:
        return str(x)


def fmt_pct(x, decimals=1):
    if x is None or x == "":
        return ""
    try:
        v = float(x)
        if v <= 1.5:
            v *= 100
        return f"{v:.{decimals}f}%"
    except Exception:
        return str(x)


def fmt_date(x):
    if x is None or x == "":
        return ""
    if isinstance(x, (datetime, date)):
        return x.strftime("%Y-%m-%d")
    return str(x)


def slide_portada(prs, data, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, GREEN_DEEP)
    # Banda lateral
    add_rect(s, 0, 0, 0.30, 7.5, AMBER)
    # Marca
    add_text(s, 0.7, 0.6, 12, 0.4, "INGENIO LA CORONA",
             font=HEAD_FONT, size=14, bold=True, color=AMBER)
    add_text(s, 0.7, 0.95, 12, 0.4, "Departamento de Tecnología de Información",
             font=HEAD_FONT, size=11, italic=True, color=AMBER_LIGHT)
    # Título grande
    add_text(s, 0.7, 2.3, 12, 1.1, "Reporte Ejecutivo de TI",
             font=HEAD_FONT, size=46, bold=True, color=WHITE)
    add_text(s, 0.7, 3.4, 12, 0.6, f"{data['inicio']['mes']} · {data['inicio']['anio']}",
             font=HEAD_FONT, size=28, bold=True, color=AMBER)
    # Línea divisoria
    add_rect(s, 0.7, 4.2, 1.8, 0.06, AMBER)
    # Pie de portada
    info = [
        ("Periodo de zafra", str(data['inicio']['zafra'])),
        ("Responsable", str(data['inicio']['responsable'])),
        ("Fecha de elaboración", str(data['inicio']['fecha'])),
    ]
    for i, (k, v) in enumerate(info):
        add_text(s, 0.7, 4.55 + i*0.45, 5, 0.35, k.upper(),
                 size=10, bold=True, color=AMBER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 3.2, 4.55 + i*0.45, 9, 0.35, v,
                 size=12, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # Ribbon inferior
    add_text(s, 0.7, 6.7, 12, 0.4, "Documento confidencial · Solo para uso de Dirección",
             size=9, italic=True, color=AMBER_LIGHT)


def slide_resumen(prs, data, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    slide_title(s, "Resumen del mes", "Mensaje clave y puntos relevantes")

    # Caja del mensaje clave
    add_rect(s, 0.5, 1.7, 12.3, 1.3, GREEN_DEEP)
    add_text(s, 0.85, 1.85, 11.7, 0.30, "MENSAJE CLAVE PARA DIRECCIÓN",
             size=10, bold=True, color=AMBER)
    add_text(s, 0.85, 2.15, 11.7, 0.85, str(data['inicio']['mensaje']),
             font=HEAD_FONT, size=18, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.TOP)

    # Puntos
    add_text(s, 0.5, 3.25, 12.3, 0.4, "Puntos relevantes",
             size=14, bold=True, color=GREEN_DEEP)
    y = 3.75
    for i, punto in enumerate(data['inicio']['puntos'][:5], 1):
        # Círculo numerado
        add_circle(s, 0.55, y, 0.55, AMBER)
        add_text(s, 0.55, y, 0.55, 0.55, str(i),
                 font=HEAD_FONT, size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 1.3, y + 0.07, 11.5, 0.45, str(punto),
                 size=13, color=DARK, anchor=MSO_ANCHOR.TOP)
        y += 0.62

    page_chrome(s, page, total, f"{data['inicio']['mes']} {data['inicio']['anio']}")


def add_native_table(slide, x, y, w, headers, rows_data, col_widths_pct,
                     row_height=0.35, header_height=0.40,
                     header_fill=GREEN_DEEP, header_text=WHITE,
                     stripe_fill=GREEN_LIGHT, header_size=10, body_size=10,
                     align_overrides=None):
    """Construye una tabla manual con rectángulos para mejor control visual."""
    align_overrides = align_overrides or {}
    n_cols = len(headers)
    col_w = [w * (p / 100.0) for p in col_widths_pct]
    # Header
    cx = x
    for i, h in enumerate(headers):
        add_rect(slide, cx, y, col_w[i], header_height, header_fill)
        align = align_overrides.get(i, PP_ALIGN.LEFT)
        add_text(slide, cx + 0.08, y, col_w[i] - 0.16, header_height, str(h),
                 size=header_size, bold=True, color=header_text,
                 align=align, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[i]
    # Rows
    cy = y + header_height
    for r_idx, row in enumerate(rows_data):
        # zebra stripe
        if r_idx % 2 == 1:
            add_rect(slide, x, cy, w, row_height, stripe_fill)
        cx = x
        for i, val in enumerate(row):
            align = align_overrides.get(i, PP_ALIGN.LEFT)
            text_color = DARK
            bold = False
            # Estado pill case: render como pill
            if isinstance(val, dict) and val.get("type") == "pill":
                pill_w = min(col_w[i] - 0.16, 1.0)
                state_pill(slide, cx + (col_w[i] - pill_w) / 2, cy + 0.05, pill_w, row_height - 0.10, val["label"])
            else:
                add_text(slide, cx + 0.08, cy, col_w[i] - 0.16, row_height, str(val),
                         size=body_size, bold=bold, color=text_color,
                         align=align, anchor=MSO_ANCHOR.MIDDLE)
            cx += col_w[i]
        cy += row_height
    return cy


def slide_disponibilidad(prs, data, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    slide_title(s, "Operación", "Disponibilidad de sistemas críticos")

    # Stat callout
    items = data["disponibilidad"]
    if items:
        avg = sum(float(r["real"] or 0) for r in items) / len(items)
        verde = sum(1 for r in items if r["estado"] == "Verde")
        amarillo = sum(1 for r in items if r["estado"] == "Amarillo")
        rojo = sum(1 for r in items if r["estado"] == "Rojo")
    else:
        avg, verde, amarillo, rojo = 0, 0, 0, 0

    # Tres callouts
    callouts = [
        (f"{avg:.2f}%", "Disponibilidad promedio", GREEN_DEEP),
        (f"{verde}", "Sistemas en verde", GREEN_OK),
        (f"{amarillo + rojo}", "Sistemas con alerta", AMBER if rojo == 0 else RED),
    ]
    for i, (val, lab, col) in enumerate(callouts):
        x = 0.5 + i * 4.30
        add_rect(s, x, 1.65, 4.10, 1.05, col)
        add_text(s, x, 1.75, 4.10, 0.55, val,
                 font=HEAD_FONT, size=32, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, 2.30, 4.10, 0.30, lab,
                 size=11, color=AMBER_LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Tabla
    rows = []
    for r in items:
        rows.append([
            r["sistema"],
            r["categoria"],
            f"{float(r['obj'] or 0):.2f}%",
            f"{float(r['real'] or 0):.2f}%",
            f"{float(r['horas'] or 0):.1f} h",
            {"type": "pill", "label": r["estado"] or "—"},
            r["comentario"] or "",
        ])
    add_native_table(
        s, 0.5, 2.95, 12.3,
        ["Sistema", "Categoría", "Objetivo", "Real", "Tiempo fuera", "Estado", "Comentario"],
        rows,
        col_widths_pct=[20, 13, 9, 9, 12, 11, 26],
        row_height=0.36, header_height=0.40,
        align_overrides={2: PP_ALIGN.CENTER, 3: PP_ALIGN.CENTER, 4: PP_ALIGN.CENTER, 5: PP_ALIGN.CENTER},
    )

    page_chrome(s, page, total, f"{data['inicio']['mes']} {data['inicio']['anio']}")


def slide_incidentes(prs, data, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    slide_title(s, "Operación", "Incidentes mayores del mes")

    items = data["incidentes"]
    n = len(items)
    h_total = sum(float(r["horas"] or 0) for r in items)
    sev_count = {"Crítica": 0, "Alta": 0, "Media": 0, "Baja": 0}
    for r in items:
        sev = r["severidad"] or "Baja"
        sev_count[sev] = sev_count.get(sev, 0) + 1

    # Callouts
    callouts = [
        (str(n), "Incidentes en el mes", GREEN_DEEP),
        (f"{h_total:.1f} h", "Tiempo acumulado de afectación", AMBER),
        (str(sev_count.get("Crítica", 0)), "Severidad Crítica", RED if sev_count.get("Crítica", 0) > 0 else GREEN_OK),
    ]
    for i, (val, lab, col) in enumerate(callouts):
        x = 0.5 + i * 4.30
        add_rect(s, x, 1.65, 4.10, 1.05, col)
        add_text(s, x, 1.75, 4.10, 0.55, val,
                 font=HEAD_FONT, size=32, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, 2.30, 4.10, 0.30, lab,
                 size=11, color=AMBER_LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    if n == 0:
        add_text(s, 0.5, 4.0, 12.3, 1.0,
                 "Sin incidentes mayores reportados durante el mes.",
                 size=18, italic=True, color=GREEN_OK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    else:
        rows = []
        for r in items:
            rows.append([
                fmt_date(r["fecha"]),
                r["sistema"] or "",
                {"type": "pill", "label": r["severidad"] or "—"},
                r["descripcion"] or "",
                f"{float(r['horas'] or 0):.1f} h",
                r["accion"] or "",
            ])
        add_native_table(
            s, 0.5, 2.95, 12.3,
            ["Fecha", "Sistema", "Severidad", "Descripción", "Tiempo", "Acción correctiva"],
            rows,
            col_widths_pct=[10, 16, 11, 30, 8, 25],
            row_height=0.65, header_height=0.40, body_size=10,
            align_overrides={2: PP_ALIGN.CENTER, 4: PP_ALIGN.CENTER},
        )

    page_chrome(s, page, total, f"{data['inicio']['mes']} {data['inicio']['anio']}")


def slide_proyectos(prs, data, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    slide_title(s, "Proyectos", "Avance de iniciativas estratégicas")

    items = data["proyectos"]
    n = len(items)
    verde = sum(1 for r in items if r["estado"] == "Verde")
    amarillo = sum(1 for r in items if r["estado"] == "Amarillo")
    rojo = sum(1 for r in items if r["estado"] == "Rojo")
    avg = sum(float(r["avance"] or 0) for r in items) / max(n, 1)
    if avg <= 1.5:
        avg_disp = avg * 100
    else:
        avg_disp = avg

    callouts = [
        (str(n), "Proyectos activos", GREEN_DEEP),
        (f"{avg_disp:.0f}%", "Avance promedio", AMBER),
        (f"{verde}/{amarillo}/{rojo}", "Verde / Amarillo / Rojo", GREEN_MID),
    ]
    for i, (val, lab, col) in enumerate(callouts):
        x = 0.5 + i * 4.30
        add_rect(s, x, 1.65, 4.10, 1.05, col)
        add_text(s, x, 1.75, 4.10, 0.55, val,
                 font=HEAD_FONT, size=30, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, 2.30, 4.10, 0.30, lab,
                 size=11, color=AMBER_LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    rows = []
    for r in items:
        adv = float(r["avance"] or 0)
        if adv <= 1.5:
            adv *= 100
        rows.append([
            r["proyecto"] or "",
            r["responsable"] or "",
            {"type": "pill", "label": r["estado"] or "—"},
            f"{adv:.0f}%",
            fmt_date(r["fin"]),
            r["hito"] or "",
        ])
    add_native_table(
        s, 0.5, 2.95, 12.3,
        ["Proyecto", "Responsable", "Estado", "Avance", "Fin previsto", "Próximo hito"],
        rows,
        col_widths_pct=[26, 16, 10, 9, 13, 26],
        row_height=0.50, header_height=0.40,
        align_overrides={2: PP_ALIGN.CENTER, 3: PP_ALIGN.CENTER, 4: PP_ALIGN.CENTER},
    )

    page_chrome(s, page, total, f"{data['inicio']['mes']} {data['inicio']['anio']}")


def slide_presupuesto(prs, data, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    slide_title(s, "Finanzas", "Presupuesto y ejecución del mes")

    items = data["presupuesto"]
    tot_p = sum(float(r["presupuesto"] or 0) for r in items)
    tot_e = sum(float(r["ejecutado"] or 0) for r in items)
    pct = (tot_e / tot_p) if tot_p else 0
    var = tot_p - tot_e

    callouts = [
        (fmt_money(tot_p), "Presupuesto del mes", GREEN_DEEP),
        (fmt_money(tot_e), "Ejecutado", AMBER),
        (f"{pct*100:.1f}%", "Ejecución", GREEN_OK if 0.85 <= pct <= 1.05 else AMBER),
    ]
    for i, (val, lab, col) in enumerate(callouts):
        x = 0.5 + i * 4.30
        add_rect(s, x, 1.65, 4.10, 1.05, col)
        add_text(s, x, 1.75, 4.10, 0.55, val,
                 font=HEAD_FONT, size=26, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, 2.30, 4.10, 0.30, lab,
                 size=11, color=AMBER_LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Tabla de presupuesto
    rows = []
    for r in items:
        p = float(r["presupuesto"] or 0)
        e = float(r["ejecutado"] or 0)
        v = p - e
        ep = (e / p) if p else 0
        rows.append([
            r["categoria"] or "",
            fmt_money(p),
            fmt_money(e),
            fmt_money(v),
            f"{ep*100:.0f}%",
        ])
    rows.append(["TOTAL", fmt_money(tot_p), fmt_money(tot_e), fmt_money(var), f"{pct*100:.0f}%"])

    end_y = add_native_table(
        s, 0.5, 2.95, 12.3,
        ["Categoría", "Presupuesto", "Ejecutado", "Variación", "% Ejec."],
        rows,
        col_widths_pct=[40, 17, 17, 14, 12],
        row_height=0.34, header_height=0.40,
        align_overrides={1: PP_ALIGN.RIGHT, 2: PP_ALIGN.RIGHT, 3: PP_ALIGN.RIGHT, 4: PP_ALIGN.CENTER},
    )

    page_chrome(s, page, total, f"{data['inicio']['mes']} {data['inicio']['anio']}")


def slide_seguridad(prs, data, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    slide_title(s, "Seguridad y cumplimiento", "Indicadores del mes")

    items = data["seguridad"]
    # Cuadrícula 2 columnas x N filas, mostrar 8 KPIs
    cols, rows_n = 4, 2
    cell_w = 3.05
    cell_h = 1.50
    start_x, start_y = 0.5, 1.65
    gap = 0.10

    items = items[: cols * rows_n]
    for idx, kpi in enumerate(items):
        col = idx % cols
        row = idx // cols
        x = start_x + col * (cell_w + gap)
        y = start_y + row * (cell_h + gap)
        # Card
        add_rect(s, x, y, cell_w, cell_h, CREAM)
        add_rect(s, x, y, cell_w, 0.08, AMBER)
        # Indicador
        ind = str(kpi["indicador"] or "")
        add_text(s, x + 0.18, y + 0.18, cell_w - 0.36, 0.45, ind,
                 size=10, bold=True, color=GRAY, anchor=MSO_ANCHOR.TOP)
        # Valor grande
        actual = kpi["actual"]
        try:
            val_disp = f"{int(float(actual)):,}"
        except Exception:
            val_disp = str(actual or "—")
        add_text(s, x + 0.18, y + 0.55, cell_w - 0.36, 0.55, val_disp,
                 font=HEAD_FONT, size=28, bold=True, color=GREEN_DEEP, anchor=MSO_ANCHOR.MIDDLE)
        # Comparación con mes anterior
        try:
            a = float(actual or 0); b = float(kpi["anterior"] or 0)
            arrow = "▲" if a > b else ("▼" if a < b else "=")
            color_arr = AMBER if a > b else (GREEN_OK if a < b else GRAY)
            comp = f"{arrow}  vs. {b:,.0f} mes anterior"
        except Exception:
            arrow = "—"; comp = ""
            color_arr = GRAY
        add_text(s, x + 0.18, y + 1.10, cell_w - 0.36, 0.30, comp,
                 size=10, color=color_arr, anchor=MSO_ANCHOR.TOP)

    page_chrome(s, page, total, f"{data['inicio']['mes']} {data['inicio']['anio']}")


def slide_licencias(prs, data, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    slide_title(s, "Licencias", "Vencimientos y costo de renovaciones próximas")

    items = data["licencias"]
    n = len(items)
    tot_costo = sum(float(r["costo"] or 0) for r in items)
    en_90d = 0
    for r in items:
        try:
            d = float(r["dias"] or 999)
            if d <= 90:
                en_90d += 1
        except Exception:
            pass

    callouts = [
        (str(n), "Licencias monitoreadas", GREEN_DEEP),
        (str(en_90d), "Vencen en ≤ 90 días", AMBER if en_90d > 0 else GREEN_OK),
        (fmt_money(tot_costo), "Costo total renovación", GREEN_MID),
    ]
    for i, (val, lab, col) in enumerate(callouts):
        x = 0.5 + i * 4.30
        add_rect(s, x, 1.65, 4.10, 1.05, col)
        add_text(s, x, 1.75, 4.10, 0.55, val,
                 font=HEAD_FONT, size=26, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, 2.30, 4.10, 0.30, lab,
                 size=11, color=AMBER_LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    rows = []
    for r in items:
        try:
            dias = int(float(r["dias"] or 0))
        except Exception:
            dias = "—"
        rows.append([
            r["producto"] or "",
            r["proveedor"] or "",
            f"{int(float(r['cantidad'] or 0)):,}" if r["cantidad"] not in (None, "") else "",
            fmt_date(r["vence"]),
            f"{dias} d" if isinstance(dias, int) else dias,
            r["estado"] or "",
            fmt_money(r["costo"]),
        ])
    add_native_table(
        s, 0.5, 2.95, 12.3,
        ["Producto", "Proveedor", "Cant.", "Vence", "Días", "Estado", "Costo"],
        rows,
        col_widths_pct=[28, 14, 8, 12, 9, 17, 12],
        row_height=0.40, header_height=0.40,
        align_overrides={2: PP_ALIGN.CENTER, 3: PP_ALIGN.CENTER, 4: PP_ALIGN.CENTER, 6: PP_ALIGN.RIGHT},
    )

    page_chrome(s, page, total, f"{data['inicio']['mes']} {data['inicio']['anio']}")


def slide_inventario(prs, data, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    slide_title(s, "Activos", "Inventario — visión ejecutiva")

    items = data["inventario"]
    tot = sum(int(float(r["cantidad"] or 0)) for r in items)
    rep_3 = sum(int(float(r["rep_3m"] or 0)) for r in items)
    rep_12 = sum(int(float(r["rep_12m"] or 0)) for r in items)

    callouts = [
        (f"{tot:,}", "Equipos en operación", GREEN_DEEP),
        (str(rep_3), "Reemplazos < 3 meses", AMBER if rep_3 > 0 else GREEN_OK),
        (str(rep_12), "Reemplazos 3-12 meses", GREEN_MID),
    ]
    for i, (val, lab, col) in enumerate(callouts):
        x = 0.5 + i * 4.30
        add_rect(s, x, 1.65, 4.10, 1.05, col)
        add_text(s, x, 1.75, 4.10, 0.55, val,
                 font=HEAD_FONT, size=28, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, 2.30, 4.10, 0.30, lab,
                 size=11, color=AMBER_LIGHT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    rows = []
    for r in items:
        rows.append([
            r["categoria"] or "",
            f"{int(float(r['cantidad'] or 0)):,}",
            f"{float(r['edad'] or 0):.1f}",
            f"{int(float(r['rep_3m'] or 0)):,}",
            f"{int(float(r['rep_12m'] or 0)):,}",
            r["comentario"] or "",
        ])
    add_native_table(
        s, 0.5, 2.95, 12.3,
        ["Categoría", "Cantidad", "Edad prom. (años)", "Reemplazo <3m", "Reemplazo 3-12m", "Comentario"],
        rows,
        col_widths_pct=[28, 11, 14, 13, 14, 20],
        row_height=0.42, header_height=0.45,
        align_overrides={1: PP_ALIGN.CENTER, 2: PP_ALIGN.CENTER, 3: PP_ALIGN.CENTER, 4: PP_ALIGN.CENTER},
    )

    page_chrome(s, page, total, f"{data['inicio']['mes']} {data['inicio']['anio']}")


def slide_cierre(prs, data, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, GREEN_DEEP)
    add_rect(s, 0, 0, 0.30, 7.5, AMBER)
    add_text(s, 0.7, 0.7, 12, 0.5, "CIERRE DEL REPORTE",
             size=12, bold=True, color=AMBER)
    add_text(s, 0.7, 1.2, 12, 0.7, "Próximos pasos y compromisos",
             font=HEAD_FONT, size=32, bold=True, color=WHITE)

    # Próximos pasos derivados automáticamente
    pasos = []
    # 1. Riesgos en proyectos
    riesgos = [r for r in data["proyectos"] if r["estado"] in ("Amarillo", "Rojo")]
    if riesgos:
        for r in riesgos[:2]:
            pasos.append(f"Atender riesgo en proyecto «{r['proyecto']}» — {r['riesgo']}")
    # 2. Licencias por vencer
    venc = []
    for r in data["licencias"]:
        try:
            if float(r["dias"] or 999) <= 90:
                venc.append(r["producto"])
        except Exception:
            pass
    if venc:
        pasos.append(f"Cerrar renovación de licencias críticas: {', '.join(venc[:3])}")
    # 3. Hallazgos abiertos (si hay)
    for r in data["seguridad"]:
        if "auditoría" in str(r["indicador"]).lower():
            try:
                if float(r["actual"] or 0) > 0:
                    pasos.append(f"Cerrar {int(float(r['actual']))} hallazgos de auditoría aún abiertos.")
            except Exception:
                pass
    # 4. Genérico
    pasos.append("Sostener disponibilidad ≥ 99.5% en sistemas críticos durante el resto de zafra.")
    pasos.append("Revisión mensual del avance presupuestal con Dirección.")
    pasos = pasos[:5]

    add_text(s, 0.7, 2.3, 12, 0.4, "Próximos pasos",
             size=14, bold=True, color=AMBER)
    for i, p in enumerate(pasos, 1):
        y = 2.7 + (i - 1) * 0.55
        add_circle(s, 0.7, y + 0.05, 0.30, AMBER)
        add_text(s, 0.7, y + 0.05, 0.30, 0.30, str(i),
                 font=HEAD_FONT, size=11, bold=True, color=GREEN_DEEP,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 1.2, y, 11.5, 0.5, p,
                 size=13, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    # Footer con responsable
    add_rect(s, 0.7, 6.5, 11.9, 0.04, AMBER)
    add_text(s, 0.7, 6.6, 12, 0.4,
             f"Elaborado por: {data['inicio']['responsable']}    ·    Fecha: {data['inicio']['fecha']}",
             size=11, italic=True, color=AMBER_LIGHT)


# ============= Pipeline =============

def main():
    src = sys.argv[1] if len(sys.argv) > 1 else "Reporte_TI_Plantilla.xlsx"
    src = Path(src)
    if not src.exists():
        print(f"No se encontró {src}")
        sys.exit(1)

    data = read_workbook(src)
    mes = str(data["inicio"]["mes"])
    anio = str(data["inicio"]["anio"])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    total = 10
    slide_portada(prs, data, total)
    slide_resumen(prs, data, 2, total)
    slide_disponibilidad(prs, data, 3, total)
    slide_incidentes(prs, data, 4, total)
    slide_proyectos(prs, data, 5, total)
    slide_presupuesto(prs, data, 6, total)
    slide_seguridad(prs, data, 7, total)
    slide_licencias(prs, data, 8, total)
    slide_inventario(prs, data, 9, total)
    slide_cierre(prs, data, 10, total)

    out = Path(f"Reporte_TI_Ejecutivo_{mes}_{anio}.pptx")
    prs.save(out)
    print(f"Generado: {out}")


if __name__ == "__main__":
    main()
